import os
import csv
import re
import pdfplumber
from pptx import Presentation
from docx import Document
import PyPDF2
import pandas as pd

def load_terms(terms_file):
    terms_df = pd.read_excel(terms_file)
    term_col = next((col for col in terms_df.columns if col.strip().lower() == 'term'), None)
    if not term_col:
        raise ValueError("No column named 'Term' found in the Excel file.")

    terms_dict = {term.strip().lower(): term.strip() for term in terms_df[term_col] if isinstance(term, str)}
    return terms_dict

def search_terms_in_text(terms_dict, text):
    matches = []
    for term_lower, term_original in terms_dict.items():
        pattern = r'\b' + re.escape(term_lower) + r'\b'
        if re.search(pattern, text, flags=re.IGNORECASE):
            matches.append(term_original)
    return matches

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text.strip():
                    text.append((slide_index, shape.text.strip()))
            elif shape.shape_type == 19:
                table = shape.table
                for row in table.rows:
                    row_text = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame]
                    if any(row_text):
                        text.append((slide_index, " | ".join(row_text)))
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text.append((1, paragraph.text.strip()))

    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells]
            if any(row_text):
                text.append((1, " | ".join(row_text)))

    return text


def extract_text_from_pdf(file_path):
    text = []

    with pdfplumber.open(file_path) as pdf:
        for page_index, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text()
            if page_text:
                for paragraph in page_text.split("\n"):
                    if paragraph.strip():
                        text.append((page_index, paragraph.strip()))

            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    row_text = [cell.strip() for cell in row if cell]
                    if any(row_text):
                        text.append((page_index, " | ".join(row_text)))

    return text

def process_file(file_path, terms_dict):
    extension = os.path.splitext(file_path)[1].lower()

    if extension == ".docx":
        text_with_pages = extract_text_from_docx(file_path)
    elif extension == ".pptx":
        text_with_pages = extract_text_from_pptx(file_path)
    elif extension == ".pdf":
        text_with_pages = extract_text_from_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file type: {extension}")

    results = []
    seen_terms = set()

    for page, paragraph in text_with_pages:
        if len(seen_terms) == len(terms_dict):
            break

        matches = search_terms_in_text(terms_dict, paragraph)
        for term in matches:
            if term not in seen_terms:
                results.append([page, term, paragraph.strip(), ""])
                seen_terms.add(term)
                break

    return results

def write_results_to_csv(results, output_file):
    with open(output_file, mode="w", newline="", encoding="utf-8") as csv_file:
        writer = csv.writer(csv_file, quoting=csv.QUOTE_ALL)
        # writer.writerow(["Page/Slide", "Term", "Paragraph", "Comments"])
        for result in results:
            result[2] = result[2].replace("\n", " ")
            writer.writerow(result)


def main():
    terms_file = "/Users/xxliu95/Documents/my_VBA_scripts/terms.xlsx"
    input_folder = "/Users/xxliu95/Documents/SAP/Liuxin"
    output_folder = "/Users/xxliu95/Documents/my_VBA_scripts/Output"

    os.makedirs(output_folder, exist_ok=True)

    terms_dict = load_terms(terms_file)
    print(f"Loaded {len(terms_dict)} terms from dictionary.")

    allowed_extensions = {".docx", ".pptx", ".pdf"}
    for file_name in os.listdir(input_folder):
        file_path = os.path.join(input_folder, file_name)
        if os.path.isfile(file_path) and os.path.splitext(file_name)[1].lower() in allowed_extensions:
            print(f"Processing {file_name}...")
            results = process_file(file_path, terms_dict)

            base_name = os.path.splitext(file_name)[0]
            output_file = os.path.join(output_folder, f"Source Analysis_{base_name}.csv")

            write_results_to_csv(results, output_file)
            print(f"Output: {output_file} \n")


if __name__ == "__main__":
    main()
